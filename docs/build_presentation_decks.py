from pathlib import Path
from datetime import date
import re

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
LOGO = ROOT / "favicon.png"
SNAPSHOT_DATE = date.today().isoformat()

PRIMARY = RGBColor(0x6D, 0x11, 0x16)
SECONDARY = RGBColor(0xC2, 0x9B, 0x6D)
BACKGROUND = RGBColor(0xF6, 0xF2, 0xEB)
TEXT = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x67, 0x67, 0x67)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def collect_inventory() -> dict:
    frontend = ROOT / "frontend"
    backend_api = ROOT / "backend" / "api"
    role_manuals = DOCS / "role_manuals"
    dfd_dir = DOCS / "DFDs"
    schema_path = ROOT / "database" / "schema.sql"
    schema_text = schema_path.read_text(encoding="utf-8", errors="ignore")
    table_count = len(re.findall(r"^CREATE TABLE IF NOT EXISTS `([^`]+)`", schema_text, re.M))
    return {
        "html_count": len(list(frontend.glob("*.html"))),
        "js_count": len(list(frontend.rglob("*.js"))),
        "php_api_count": len(list(backend_api.glob("*.php"))),
        "table_count": table_count,
        "role_manual_count": len(list(role_manuals.glob("PensionApp_Role_Manual_*.md"))),
        "dfd_count": len([path for path in dfd_dir.glob("*.md") if path.name.lower() != "readme.md"]),
    }


def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND


def add_logo(slide, left, top, width):
    if LOGO.exists() and LOGO.stat().st_size > 0:
        slide.shapes.add_picture(str(LOGO), left, top, width=width)


def add_header_footer(slide, section_label: str, slide_no: int):
    top_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.55))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = PRIMARY
    top_band.line.fill.background()

    bottom_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(7.1), SLIDE_W, Inches(0.4))
    bottom_band.fill.solid()
    bottom_band.fill.fore_color.rgb = PRIMARY
    bottom_band.line.fill.background()

    add_logo(slide, Inches(0.28), Inches(0.08), Inches(0.34))

    header_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.1), Inches(8.4), Inches(0.28))
    p = header_box.text_frame.paragraphs[0]
    p.text = section_label
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE

    footer_left = slide.shapes.add_textbox(Inches(0.35), Inches(7.16), Inches(5.8), Inches(0.18))
    p = footer_left.text_frame.paragraphs[0]
    p.text = "UPS PensionsGo"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE

    footer_right = slide.shapes.add_textbox(Inches(12.2), Inches(7.16), Inches(0.7), Inches(0.18))
    p = footer_right.text_frame.paragraphs[0]
    p.text = str(slide_no)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.92), Inches(11.8), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.58), Inches(11.0), Inches(0.45))
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12.5)
        p.font.color.rgb = MUTED


def add_bullets(slide, bullets, left=Inches(0.9), top=Inches(2.0), width=Inches(11.6), height=Inches(4.6), font_size=Pt(18)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = bullet
        p.level = 0
        p.font.size = font_size
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
        p.bullet = True


def add_two_column_bullets(slide, left_title, left_items, right_title, right_items, item_font_size=Pt(15.5)):
    card_w = Inches(5.8)
    for idx, (title, items, left) in enumerate([
        (left_title, left_items, Inches(0.75)),
        (right_title, right_items, Inches(6.55)),
    ]):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.95), card_w, Inches(4.85))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = SECONDARY

        title_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.12), card_w - Inches(0.4), Inches(0.35))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        add_bullets(slide, items, left + Inches(0.18), Inches(2.55), card_w - Inches(0.35), Inches(4.0), font_size=item_font_size)


def add_metrics_row(slide, metrics):
    card_w = Inches(2.95)
    gap = Inches(0.2)
    left = Inches(0.78)
    for value, label, detail in metrics:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.0), card_w, Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = SECONDARY

        value_box = slide.shapes.add_textbox(left + Inches(0.18), Inches(2.2), card_w - Inches(0.35), Inches(0.58))
        p = value_box.text_frame.paragraphs[0]
        p.text = value
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER

        label_box = slide.shapes.add_textbox(left + Inches(0.18), Inches(2.85), card_w - Inches(0.35), Inches(0.38))
        p = label_box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT
        p.alignment = PP_ALIGN.CENTER

        detail_box = slide.shapes.add_textbox(left + Inches(0.18), Inches(3.28), card_w - Inches(0.35), Inches(0.65))
        p = detail_box.text_frame.paragraphs[0]
        p.text = detail
        p.font.size = Pt(10.5)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.CENTER

        left += card_w + gap


def add_process_flow(slide, steps):
    width = Inches(2.35)
    gap = Inches(0.18)
    left = Inches(0.62)
    for idx, step in enumerate(steps):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.35), width, Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = SECONDARY

        num_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.88), Inches(2.08), Inches(0.55), Inches(0.55))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = PRIMARY
        num_box.line.fill.background()
        p = num_box.text_frame.paragraphs[0]
        p.text = str(idx + 1)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        title_box = slide.shapes.add_textbox(left + Inches(0.15), Inches(2.8), width - Inches(0.3), Inches(0.42))
        p = title_box.text_frame.paragraphs[0]
        p.text = step["title"]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        body_box = slide.shapes.add_textbox(left + Inches(0.18), Inches(3.28), width - Inches(0.35), Inches(1.0))
        p = body_box.text_frame.paragraphs[0]
        p.text = step["body"]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT

        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left + width + Inches(0.03), Inches(3.05), Inches(0.25), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PRIMARY
            arrow.line.fill.background()

        left += width + gap


def add_title_slide(prs, title, subtitle, audience_tag, date_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.7))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY
    accent.line.fill.background()

    add_logo(slide, Inches(0.6), Inches(1.0), Inches(1.35))

    title_box = slide.shapes.add_textbox(Inches(2.2), Inches(1.08), Inches(9.9), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    subtitle_box = slide.shapes.add_textbox(Inches(2.24), Inches(2.0), Inches(9.5), Inches(0.55))
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT

    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.22), Inches(2.75), Inches(3.1), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = SECONDARY
    badge.line.fill.background()
    p = badge.text_frame.paragraphs[0]
    p.text = audience_tag
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    meta_box = slide.shapes.add_textbox(Inches(2.24), Inches(4.6), Inches(5.8), Inches(1.0))
    p = meta_box.text_frame.paragraphs[0]
    p.text = "UPS PensionsGo"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = meta_box.text_frame.add_paragraph()
    p.text = "Prepared from the live repository and documentation set"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT
    p = meta_box.text_frame.add_paragraph()
    p.text = date_text
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED

    return slide


def add_bullet_slide(prs, section, slide_no, title, subtitle, bullets, font_size=Pt(18), top=Inches(2.0), height=Inches(4.6)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, section, slide_no)
    add_title(slide, title, subtitle)
    add_bullets(slide, bullets, top=top, height=height, font_size=font_size)
    return slide


def add_two_column_slide(prs, section, slide_no, title, subtitle, left_title, left_items, right_title, right_items, item_font_size=Pt(15.5)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, section, slide_no)
    add_title(slide, title, subtitle)
    add_two_column_bullets(slide, left_title, left_items, right_title, right_items, item_font_size=item_font_size)
    return slide


def add_metrics_slide(prs, section, slide_no, title, subtitle, metrics, bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, section, slide_no)
    add_title(slide, title, subtitle)
    add_metrics_row(slide, metrics)
    if bullets:
        add_bullets(slide, bullets, left=Inches(0.95), top=Inches(4.55), width=Inches(11.4), height=Inches(1.7), font_size=Pt(14.5))
    return slide


def add_process_slide(prs, section, slide_no, title, subtitle, steps, closing_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, section, slide_no)
    add_title(slide, title, subtitle)
    add_process_flow(slide, steps)
    if closing_bullets:
        add_bullets(slide, closing_bullets, left=Inches(0.9), top=Inches(5.0), width=Inches(11.6), height=Inches(1.4), font_size=Pt(13.5))
    return slide


def build_technical_deck():
    inventory = collect_inventory()
    prs = make_presentation()
    add_title_slide(
        prs,
        "PensionApp Technical Team Briefing",
        "Architecture, domain model, governance, and engineering priorities",
        "Audience: Technical Team",
        f"Repository snapshot: {SNAPSHOT_DATE}",
    )
    add_bullet_slide(
        prs,
        "Technical Overview",
        2,
        "Why This Platform Matters",
        "What the technical team needs to hold in view",
        [
            "UPS PensionsGo is a full-stack pension administration platform with workflow, registry, payroll, arrears, messaging, feedback, and pensioner self-service.",
            f"The current codebase spans {inventory['html_count']} HTML pages, {inventory['js_count']} JavaScript files, {inventory['php_api_count']} PHP API files, and {inventory['table_count']} database tables.",
            "The shared runtime spine is concentrated in backend/config.php, which combines settings, security, schema hardening, and reusable domain helpers.",
            "The strongest business-data thread runs from tb_staffdue to tb_fileregistry and then into tasks, payroll, arrears, documents, and pensioner services.",
        ],
    )
    add_metrics_slide(
        prs,
        "Technical Overview",
        3,
        "System Footprint",
        "Current repository and data-model scale",
        [
            (str(inventory["html_count"]), "HTML pages", "Operational, admin, public, and pensioner interfaces"),
            (str(inventory["js_count"]), "JS files", "Vanilla JS controllers plus shared modules"),
            (str(inventory["php_api_count"]), "PHP APIs", "Small endpoint-style service layer"),
            (str(inventory["table_count"]), "DB tables", "Workflow, registry, payroll, arrears, governance, and content"),
        ],
        [
            "Runtime schema hardening remains concentrated in backend/config.php and helper files, so release discipline still matters.",
            "The ERD pack is mature enough to support onboarding, refactoring, and cross-domain debugging, including system health diagnostics coverage.",
        ],
    )
    add_two_column_slide(
        prs,
        "Architecture",
        4,
        "Architecture Stack",
        "How the application is layered today",
        "Application Layers",
        [
            "Frontend: static HTML + CSS + vanilla JS",
            "PWA shell: manifest + service worker",
            "API layer: focused PHP endpoints under backend/api",
            "Shared runtime: backend/config.php and helper files",
            "Persistence: MySQL/MariaDB schema under database/schema.sql",
            "Worker layer: notification queue CLI process",
        ],
        "Operational Characteristics",
        [
            "Role-aware navigation and API enforcement",
            "Settings-driven behavior across security, retention, notifications, and portal access",
            "RegNo-centered cross-domain joins",
            "Upload-heavy workflows for registry, payroll, messaging, and accountability",
            "Governance surfaces for delete, restore, export, backup, and cleanup",
        ],
    )
    add_two_column_slide(
        prs,
        "Domains",
        5,
        "Capability Map",
        "Nine functional domains dominate the codebase",
        "Operational Domains",
        [
            "Workflow & Tasks",
            "Claims & Intake",
            "Registry & Document Control",
            "Payroll & Suspensions",
            "Arrears & Budget",
        ],
        "Service & Governance Domains",
        [
            "Messaging & Notifications",
            "Feedback / FAQ / Terms / Podcast",
            "Pensioner Self-Service",
            "Administration & Operations",
            "ERD / documentation support assets",
        ],
    )
    add_process_slide(
        prs,
        "Data and Flow",
        6,
        "End-to-End Operating Flow",
        "The sequence the codebase is built around",
        [
            {"title": "Intake", "body": "Staff due records captured in tb_staffdue"},
            {"title": "Workflow", "body": "Queueing, status tracking, tasks, and approvals"},
            {"title": "Registry", "body": "Canonical pension file created and maintained"},
            {"title": "Payroll", "body": "Monthly cycles uploaded and matched"},
            {"title": "Accountability", "body": "Arrears, documents, exports, and governance"},
        ],
        [
            "Use this flow to reason about ownership boundaries, bug locations, and change impact.",
            "Most debugging sessions become easier once the affected stage is identified correctly.",
        ],
    )
    add_two_column_slide(
        prs,
        "Data and Flow",
        7,
        "Data Model Spine",
        "The tables that connect the grouped ERDs back into a whole",
        "Anchor Tables",
        [
            "tb_users: identity and actor trail",
            "tb_staffdue: pre-registry intake record",
            "tb_fileregistry: canonical pension file",
            "tb_tasks: operational orchestration",
            "tb_payroll_upload_cycles: monthly reconciliation spine",
            "tb_arrears_payments: financial accountability spine",
        ],
        "ERD Navigation Strategy",
        [
            "Start with docs/ERD.pdf or docs/erd.md",
            "Deep-dive via docs/erd-domains/*.md",
            "Reconnect through interdomain_links and interdomain_matrix",
            "Use _missing_tables.txt as the completeness check",
        ],
    )
    add_two_column_slide(
        prs,
        "Security and Governance",
        8,
        "Security Model",
        "Controls already present in code",
        "Core Controls",
        [
            "Role-based access with dynamic roles and permission overrides",
            "Effective-role resolution for cloned and alias roles",
            "CSRF and origin validation support",
            "Session and device-token controls",
            "Admin re-authentication and inactivity handling",
        ],
        "Governance Controls",
        [
            "Delete queues and recycle-bin restore flow",
            "Settings-driven retention and cleanup",
            "Audit logs, system logs, workflow logs, delegation logs",
            "Server-side mirrors for front-end restrictions",
            "Restricted Data Management access for admin and OC/deputy OC roles",
        ],
    )
    add_two_column_slide(
        prs,
        "Operations",
        9,
        "Deployment and Runtime Ops",
        "Where the platform expects operational discipline",
        "Runtime Model",
        [
            "PHP/Apache hosting with MySQL/MariaDB",
            "XAMPP-local and cPanel-style deployment examples",
            "SMTP or PHP mail transport",
            "CLI notification worker for queue processing",
            "Uploads stored under backend/uploads",
        ],
        "Operational Jobs",
        [
            "Notification queue processing",
            "Analytics digests and message-storage snapshots",
            "Data import/export routines",
            "Backup and restore",
            "Data cleanup and storage reporting",
        ],
    )
    add_two_column_slide(
        prs,
        "Engineering Posture",
        10,
        "Strengths and Risks",
        "Where the codebase is strong and where it needs care",
        "Current Strengths",
        [
            "Broad functional coverage in one platform",
            "Good governance instincts in access and lifecycle flows",
            "Existing ERD/documentation pack reduces onboarding cost",
            "Settings-driven behavior makes operations adaptable",
            "PWA shell improves resilience and reach",
        ],
        "Current Risks",
        [
            "Monolithic backend/config.php increases coupling",
            "Request-time schema mutation hides migration discipline",
            "Large endpoint count makes discovery hard without catalogues",
            "Automated test coverage is very limited",
            "Legacy naming overlap and historical artifacts need pruning",
        ],
    )
    add_bullet_slide(
        prs,
        "Engineering Posture",
        11,
        "Recent Repo Hygiene",
        "Current cleanup item to keep in mind",
        [
            "The overlapping fetch_staff* duplication has been reduced; only backend/api/fetch_staffdue.php remains in the repository.",
            "This is a good direction: route clarity matters in a PHP endpoint-heavy codebase.",
            "The next step is to keep naming conventions consistent so page controllers, APIs, and docs line up cleanly.",
        ],
    )
    add_bullet_slide(
        prs,
        "Roadmap",
        12,
        "Recommended Near-Term Priorities",
        "What the technical team should do next",
        [
            "Introduce formal versioned migrations so runtime ensure-functions become defensive rather than primary schema drivers.",
            "Build regression coverage around auth, registry update paths, payroll uploads, recycle-bin restore, task routing, and pensioner restrictions.",
            "Publish an API catalogue or OpenAPI baseline for the most critical operational endpoints.",
            "Split backend/config.php into coherent service modules for settings, RBAC, logs, notifications, and domain-specific utilities.",
            "Continue removing ambiguous or duplicate routes and keep docs aligned with those changes.",
        ],
    )
    add_bullet_slide(
        prs,
        "Close",
        13,
        "Technical Discussion Prompts",
        "Suggested framing for the review session",
        [
            "Which domain should be modularized first without breaking momentum?",
            "Which workflows most urgently need automated tests?",
            "Where should formal migrations begin: roles/settings, registry, or payroll?",
            "Which access-control rules deserve a central policy map next?",
        ],
    )
    return prs


def build_user_deck():
    inventory = collect_inventory()
    prs = make_presentation()
    add_title_slide(
        prs,
        "PensionApp User Briefing",
        "How to use UPS PensionsGo safely, consistently, and confidently",
        "Audience: Operational Users",
        f"Repository snapshot: {SNAPSHOT_DATE}",
    )
    add_bullet_slide(
        prs,
        "User Overview",
        2,
        "What PensionsGo Is For",
        "The platform's practical purpose for daily work",
        [
            "PensionsGo brings pension workflow, registry, claims, payroll visibility, tasks, messaging, and support tools into one place.",
            f"The current platform baseline includes {inventory['html_count']} pages, {inventory['php_api_count']} service endpoints, and a maintained schema with {inventory['table_count']} tables.",
            "It helps teams move work forward with clearer ownership, stronger traceability, and less dependence on informal follow-up.",
            "The system is role-aware, so not every user sees the same menu or actions.",
            "The goal is not only speed, but safe and accountable processing.",
        ],
    )
    add_two_column_slide(
        prs,
        "User Overview",
        3,
        "Who Uses the System",
        "Access is based on role and responsibility",
        "Main User Groups",
        [
            "Clerks and intake staff",
            "OC/Pension and deputy supervisors",
            "Workflow officers",
            "Registry and records staff",
            "Payroll and finance staff",
        ],
        "Special Access Groups",
        [
            "Administrators",
            "Support staff",
            "Pensioner users",
            "Public feedback users",
            "Content and communications managers",
        ],
    )
    add_bullet_slide(
        prs,
        "Getting Started",
        4,
        "Using the Platform Day to Day",
        "The basic pattern most users should follow",
        [
            "Log in with the correct account and wait for the page to finish loading before acting.",
            "Start from the dashboard when you need a summary, and from the module page when you already know the task.",
            "Apply filters before exporting, printing, or reporting figures.",
            "Use the system's forms, queues, and workflow buttons instead of informal workarounds.",
        ],
    )
    add_process_slide(
        prs,
        "Workflow",
        5,
        "Standard Working Flow",
        "A simple way to understand where your work fits",
        [
            {"title": "Capture", "body": "Staff due or application data is entered"},
            {"title": "Route", "body": "Tasks, statuses, and approvals move the work"},
            {"title": "Record", "body": "Registry and documents become the source of truth"},
            {"title": "Reconcile", "body": "Payroll and arrears data are checked"},
            {"title": "Support", "body": "Users message, export, review, and follow up"},
        ],
        [
            "If a case seems stuck, first work out which stage it is in before escalating.",
        ],
    )
    add_two_column_slide(
        prs,
        "Modules",
        6,
        "Core Working Modules",
        "Where users go for the main kinds of work",
        "Operational Modules",
        [
            "Dashboard",
            "Staff Due for Retirement",
            "Pension File Registry",
            "Claims & Arrears",
            "Tasks",
            "Messages",
        ],
        "Control and Support Modules",
        [
            "Life Certificate management",
            "Payroll upload history",
            "Users management",
            "Feedback management",
            "Settings",
            "Data Management",
        ],
    )
    add_two_column_slide(
        prs,
        "Modules",
        7,
        "What Good Use Looks Like",
        "Practical operating habits that keep the platform reliable",
        "Do This",
        [
            "Confirm you are updating the correct record",
            "Use the correct module for the task",
            "Keep comments and notes factual",
            "Review filters before exporting",
            "Escalate through the proper supervisor route",
        ],
        "Avoid This",
        [
            "Skipping workflow stages without instruction",
            "Sharing accounts or passwords",
            "Treating message threads as informal, off-record chat",
            "Deleting or purging data outside approved flows",
            "Changing settings casually without documenting why",
        ],
    )
    add_bullet_slide(
        prs,
        "Governance",
        8,
        "Data Management and Restore Rules",
        "Users should understand these controls even if they cannot perform them",
        [
            "Registry deletions do not simply disappear; they move through a governed delete and recycle-bin flow.",
            "Data Management access is restricted to administrators and OC/Pension-equivalent leadership roles.",
            "Staff-due deletion follows a separate approval queue, so users should use the proper request path.",
            "Restore and purge actions must be treated as controlled governance actions, not convenience shortcuts.",
        ],
    )
    add_bullet_slide(
        prs,
        "Pensioner Support",
        9,
        "Supporting Pensioner-Facing Work",
        "What staff should know when helping pensioners",
        [
            "Pensioner portal access may be enabled or disabled centrally.",
            "Pensioner lookup visibility depends on consent settings.",
            "Some contact fields are intentionally restricted in special lifecycle states.",
            "If a pensioner cannot change details, check whether the record is in a protected condition before assuming the system is broken.",
        ],
    )
    add_two_column_slide(
        prs,
        "Security",
        10,
        "Security and Good Practice",
        "Simple habits that prevent avoidable incidents",
        "Every User Should",
        [
            "Protect credentials",
            "Respect timeout and re-auth prompts",
            "Handle exports and attachments as sensitive data",
            "Use only authorized actions for the current role",
            "Report unusual behavior quickly",
        ],
        "Supervisors and Admins Should",
        [
            "Review access regularly",
            "Watch delete, restore, export, and cleanup activity carefully",
            "Use the dashboard as a prompt, then verify in the source module",
            "Document important settings changes",
            "Guide users toward the governed path, not workarounds",
        ],
    )
    add_bullet_slide(
        prs,
        "Troubleshooting",
        11,
        "When Something Goes Wrong",
        "A simple first-response checklist",
        [
            "If login fails: check credentials, role, lockout, and maintenance-mode conditions.",
            "If an action is denied: confirm both role and permission scope before escalating.",
            "If numbers look wrong: refresh, confirm filters, and verify the source module.",
            "If a deleted record is needed: check the delete queue or recycle bin instead of assuming permanent loss.",
            "If a pensioner cannot edit details: confirm whether lifecycle restrictions are in effect.",
        ],
    )
    add_bullet_slide(
        prs,
        "Close",
        12,
        "Key Message for Users",
        "The platform works best when the process is followed as designed",
        [
            "Use the right module for the right job.",
            "Let the workflow, queues, and approvals do the coordination work.",
            "Treat every change as part of an auditable operational record.",
            "Ask early when access, data state, or workflow ownership is unclear.",
        ],
    )
    return prs


def build_training_deck():
    inventory = collect_inventory()
    prs = make_presentation()
    add_title_slide(
        prs,
        "UPS PensionsGo Comprehensive Training Deck",
        "Full stakeholder training on functionality, workflows, governance, and core use cases",
        "Audience: All Stakeholders",
        f"Repository snapshot: {SNAPSHOT_DATE}",
    )

    slide_no = 2

    def next_slide_no():
        nonlocal slide_no
        value = slide_no
        slide_no += 1
        return value

    add_bullet_slide(
        prs,
        "Training Overview",
        next_slide_no(),
        "Training Objectives",
        "What this training session is designed to achieve",
        [
            "Explain what UPS PensionsGo covers across intake, workflow, registry, claims, payroll, messaging, pensioner service, and administration.",
            "Show how different stakeholder groups use the platform without assuming every user sees the same pages or actions.",
            "Walk through the end-to-end pension case lifecycle and the governance controls that protect it.",
            "Connect platform features to real user use cases, escalation paths, and support responsibilities.",
            "Point every audience group to the correct follow-on manuals, diagrams, and operational references.",
        ],
        font_size=Pt(16.5),
    )
    add_two_column_slide(
        prs,
        "Training Overview",
        next_slide_no(),
        "Who This Deck Serves",
        "The training is designed for both decision-makers and hands-on users",
        "Decision and Oversight Stakeholders",
        [
            "Project sponsors and institutional leadership",
            "Head of Pensions and line managers",
            "OC/Pension and Deputy OC/Pension roles",
            "Administrators, governance, and support leads",
            "Audit and approval-stage decision makers",
        ],
        "Operational and Service Stakeholders",
        [
            "Clerks and intake teams",
            "Writeup, file-creation, and data-entry officers",
            "Assessors, auditors, and approvers",
            "Payroll, finance, registry, and records staff",
            "Pensioners, support teams, and public-information users",
        ],
        item_font_size=Pt(14.6),
    )
    add_metrics_slide(
        prs,
        "Platform Scope",
        next_slide_no(),
        "Current Platform Footprint",
        "The training deck is based on the live repository and current documentation suite",
        [
            (str(inventory["html_count"]), "HTML pages", "Operational, admin, pensioner, and public interfaces"),
            (str(inventory["php_api_count"]), "PHP APIs", "Endpoint-style service layer for platform actions"),
            (str(inventory["table_count"]), "DB tables", "Maintained schema baseline for live operations"),
            (str(inventory["role_manual_count"]), "Role manuals", "Role-specific Word/PDF manuals already generated"),
        ],
        [
            f"The documentation pack also includes {inventory['dfd_count']} maintained DFD diagrams, a full ERD reference, and current system and user manuals.",
            "Training should stay role-aware because visibility in the UI is narrower than total system capability.",
        ],
    )
    add_two_column_slide(
        prs,
        "Platform Scope",
        next_slide_no(),
        "Capability Map",
        "The app is broader than a single registry or claims page",
        "Operational Domains",
        [
            "Staff-due intake and application handling",
            "Workflow, tasks, delegation, and status control",
            "Pension file registry and document management",
            "File movement and life-certificate compliance",
            "Claims, arrears, and accountability handling",
        ],
        "Service and Control Domains",
        [
            "Payroll uploads, suspensions, and gratuity schedules",
            "Budgeting, reports, and exports",
            "Messages, broadcasts, and notifications",
            "Feedback, FAQ, Terms, and podcast content",
            "Pensioner self-service and administration",
        ],
        item_font_size=Pt(14.4),
    )
    add_two_column_slide(
        prs,
        "Roles and Access",
        next_slide_no(),
        "Role Model and Access Principles",
        "The same system serves many roles, but not every role sees the same controls",
        "Workflow and Line Roles",
        [
            "Clerk",
            "Writeup Officer",
            "File Creator",
            "Data Entrant",
            "Assessor, Auditor, and Approver",
        ],
        "Supervision and Service Roles",
        [
            "Administrator",
            "OC/Pension and Deputy OC/Pension",
            "General internal User",
            "Pensioner",
            "Public guidance and support users",
        ],
        item_font_size=Pt(14.8),
    )
    add_process_slide(
        prs,
        "Core Workflow",
        next_slide_no(),
        "End-to-End Pension Case Lifecycle",
        "The application is designed around this operating sequence",
        [
            {"title": "Capture", "body": "Staff due and intake records are created and verified"},
            {"title": "Route", "body": "Tasks, comments, alerts, and workflow actions move the case"},
            {"title": "Record", "body": "Registry, documents, and file custody become authoritative"},
            {"title": "Reconcile", "body": "Claims, payroll, suspensions, and accountability are checked"},
            {"title": "Serve", "body": "Decision, support, pensioner service, and governance continue"},
        ],
        [
            "Each stage has distinct owners, but the audit trail and source data must remain connected from end to end.",
            "When a case seems stuck, identify the stage first before escalating or reworking it.",
        ],
    )
    add_bullet_slide(
        prs,
        "Core Workflow",
        next_slide_no(),
        "Daily Navigation Pattern",
        "A consistent usage pattern reduces errors and rework",
        [
            "Sign in with the correct account and let the platform complete session checks and role-based redirect first.",
            "Start from the dashboard when you need the operational picture, and from the module page when you already know the task.",
            "Search and filter before editing, exporting, printing, or escalating what you see.",
            "Open the detail panel, modal, or edit workspace for the exact record before changing anything.",
            "Use the governed workflow action and wait for confirmation rather than refreshing away or improvising a shortcut.",
        ],
        font_size=Pt(16.0),
    )
    add_bullet_slide(
        prs,
        "Modules",
        next_slide_no(),
        "Dashboard and Operational Oversight",
        "The dashboard is the shared command surface for most internal stakeholders",
        [
            "The dashboard brings together claims exposure, pensioner demographics, life-certificate compliance, payroll coverage, staff-due pipeline, file status, workflow performance, and feedback management.",
            "Supervisors and leadership should use it to identify pressure points, then confirm those signals in the source module before deciding.",
            "Export and reporting actions should only be used after filters and date scope are checked carefully.",
            "Some governed management views are deliberately restricted to administrators and OC/Pension-equivalent leadership roles.",
            "As a training rule: treat the dashboard as the control picture, not as the only source of record truth.",
        ],
        font_size=Pt(15.6),
    )
    add_bullet_slide(
        prs,
        "Modules",
        next_slide_no(),
        "Staff Due and Application Intake",
        "The intake workspace is the foundation for later workflow quality",
        [
            "Use Staff Due for Retirement together with Add Staff, Edit Staff Record, and View Staff to create, review, and refine retirement-bound staff records.",
            "Capture identity, service, retirement, and contact data carefully because later roles inherit these values.",
            "Use workflow actions and status updates to move work forward instead of leaving hidden assumptions in comments.",
            "Request deletion through the governed staff-due delete flow when a record should not continue in normal use.",
            "Clerks, data-entry staff, and writeup roles are the main training audiences for this area, with supervision over exceptions and queues.",
        ],
        font_size=Pt(15.5),
    )
    add_bullet_slide(
        prs,
        "Modules",
        next_slide_no(),
        "Tasks, Delegation, and Workflow Control",
        "The tasks workspace is the execution layer for routed work",
        [
            "Use My Tasks to review assigned work, add factual comments, update status, and process completion or alert queues where permitted.",
            "Delegation should follow role boundaries so the next user receives work that matches their responsibility and authority.",
            "Workflow comments should explain real progress, blockers, or evidence gaps rather than becoming side-channel conversation.",
            "Application status and task history together tell the strongest story of where a case really stands.",
            "OC/Pension, Deputy OC/Pension, administrators, and line users all need to understand this workspace, even if they use it differently.",
        ],
        font_size=Pt(15.5),
    )
    add_bullet_slide(
        prs,
        "Modules",
        next_slide_no(),
        "Pension File Registry and Document Control",
        "The registry is the authoritative post-intake source of pension record truth",
        [
            "Use the pension file registry to search the canonical file, review benefits snapshot data, inspect document index, and maintain approved fields where the role permits it.",
            "Document upload and document-viewer flows keep evidence linked to the correct record instead of scattering files across informal channels.",
            "Delete requests, recycle-bin restore, and governed purge functions protect the registry from untraceable loss.",
            "Registry work should always begin by confirming the correct file number, pensioner identity, and current workflow context.",
            "Registry quality matters to almost every stakeholder group because payroll, claims, pensioner service, and workflow all depend on it.",
        ],
        font_size=Pt(15.2),
    )
    add_two_column_slide(
        prs,
        "Modules",
        next_slide_no(),
        "File Movement and Life Certificate Control",
        "These two record-management areas are especially sensitive operationally",
        "File Movement",
        [
            "Record every movement when custody changes",
            "Capture destination, reason, and expected return detail",
            "Use return actions to close the custody loop",
            "Keep file location aligned with task and workflow state",
            "Treat movement logs as accountability records",
        ],
        "Life Certificates",
        [
            "Work against the correct year before saving",
            "Treat Submitted, Not Submitted, and Exempt as distinct states",
            "Verify the correct profile before updating contact fields",
            "Use life-certificate status for compliance follow-up and support",
            "Escalate uncertainty rather than saving a guessed status",
        ],
        item_font_size=Pt(14.3),
    )
    add_bullet_slide(
        prs,
        "Modules",
        next_slide_no(),
        "Claims, Arrears, and Accountability",
        "The claims workspace is a full operational area, not just a simple table",
        [
            "Current functionality includes arrears-ledger review, recent payments, bulk payment upload, suspension history, gratuity schedule analysis, accountability handling, and governed exports.",
            "Users must verify beneficiary identity, claim type, period, expected amount, paid amount, and remaining balance before posting or approving financial changes.",
            "Strategic views such as estate-expiry and full-pension-due reporting support supervisory planning for privileged roles.",
            "Accountability submissions and supporting files should be completed promptly where the payment context requires them.",
            "Training for this area should separate view-only audiences from users who can actually manage arrears entries or suspension uploads.",
        ],
        font_size=Pt(15.0),
    )
    add_two_column_slide(
        prs,
        "Modules",
        next_slide_no(),
        "Payroll, Suspensions, Budgeting, and Reports",
        "These functions connect financial evidence, reconciliation, and management reporting",
        "Payroll and Suspensions",
        [
            "Upload payroll cycles and payment-register evidence",
            "Review matched and unmatched rows after every upload",
            "Upload suspension records only from approved source files",
            "Preserve review exports and exception outputs",
            "Leave cycle replacement or deletion to admin-authorized users",
        ],
        "Budgeting and Reports",
        [
            "Use budgeting to manage arrears and pension forecast views",
            "Use reports and exports for review, audit, and management packs",
            "Always check date, quarter, and financial-year scope first",
            "Treat exported files as controlled operational records",
            "Escalate unexplained financial anomalies quickly",
        ],
        item_font_size=Pt(14.0),
    )
    add_bullet_slide(
        prs,
        "Modules",
        next_slide_no(),
        "Messages, Broadcasts, and Notifications",
        "Collaboration is built into the platform and should remain professional and auditable",
        [
            "The messages module supports direct messaging, recipient control, attachments, unread tracking, and attachment viewing.",
            "Broadcasts and queued notifications support wider communication without requiring informal parallel channels.",
            "Message threads should be used for controlled clarifications and workflow follow-up, not casual off-record discussion.",
            "Attachments, message content, and generated notifications should be handled as sensitive operational information.",
            "Support and administrator teams should know how notification queues, digests, and storage snapshots affect communication reliability.",
        ],
        font_size=Pt(15.5),
    )
    add_bullet_slide(
        prs,
        "Pensioner Services",
        next_slide_no(),
        "Pensioner Dashboard and Consent-Based Lookup",
        "Pensioner-facing functionality is rich, but intentionally governed",
        [
            "The pensioner dashboard organizes information into profile, application, benefits, compliance, claims, and lifecycle or document views.",
            "Claims, documents, and pensioner lookup can be enabled or disabled centrally, so support teams should check settings before treating them as missing features.",
            "Profile edits are limited to approved self-service fields, and some fields become restricted in sensitive lifecycle states.",
            "Pensioner lookup is directory-style and consent-based: only pensioners who choose visibility should appear to others.",
            "Pensioner support training should focus on linkage checks, consent rules, and clear escalation to the pensions office for formal corrections.",
        ],
        font_size=Pt(15.1),
    )
    add_two_column_slide(
        prs,
        "Public and Support",
        next_slide_no(),
        "Feedback, Guidance, and Public-Facing Information",
        "The app serves both internal operations and controlled public support",
        "Guidance and Content Surfaces",
        [
            "About page for product positioning and service context",
            "FAQ page for common workflow and pension questions",
            "Terms page for governed platform-use expectations",
            "Podcast pages for public and internal explanatory videos",
            "Offline and PWA support for resilience and distribution",
        ],
        "Feedback and Service Signals",
        [
            "Feedback can be opened to public, staff, and pensioners by settings",
            "Feedback workflow supports assignment, review, resolution, closure, and export",
            "Dashboard analytics highlight overdue or unresolved submissions",
            "Content managers and admins maintain FAQ, terms, and podcast entries",
            "Support teams should use these surfaces as formal guidance channels",
        ],
        item_font_size=Pt(14.0),
    )
    add_bullet_slide(
        prs,
        "Administration",
        next_slide_no(),
        "Settings and Administration Workspace",
        "Administrators maintain the platform and supervise its operating conditions",
        [
            "The Settings workspace includes Dashboard Overview, User Management, App Settings, Security Settings, Access Control, role governance, notifications, storage, logs, audit trail, analysis, and system health.",
            "Administrators manage users, roles, permission overrides, content settings, notification behavior, and operational diagnostics from one governed workspace.",
            "Data-management actions such as backup, restore, import, export, and cleanup should be treated as controlled maintenance, not convenience shortcuts.",
            "Admin access is broad but should still respect business workflow ownership and leave a clear audit trail for exceptional interventions.",
            "Support, governance, and technical stakeholders should all understand what lives in this workspace, even when they do not operate it daily.",
        ],
        font_size=Pt(15.0),
    )
    add_two_column_slide(
        prs,
        "Governance",
        next_slide_no(),
        "Security, Sessions, Roles, and Permission Governance",
        "The platform protects data through layered controls rather than one single feature",
        "Security and Session Controls",
        [
            "Protected pages require a valid session",
            "Timeouts, warnings, and concurrent-session rules are enforced",
            "Sensitive admin actions may require re-authentication",
            "Maintenance mode can restrict ordinary users while allowing admin access",
            "Pensioner login can be centrally enabled or disabled",
        ],
        "Access and Governance Controls",
        [
            "Role-aware navigation and server-side permission checks work together",
            "Role and user permission overrides refine the default access model",
            "Delete queues and recycle-bin flows protect record integrity",
            "Audit, user, workflow, and system logs preserve evidence trails",
            "Data Management is restricted to higher-governance roles",
        ],
        item_font_size=Pt(14.0),
    )
    add_bullet_slide(
        prs,
        "Governance",
        next_slide_no(),
        "Data Management, Export, Backup, Restore, and Cleanup",
        "Every stakeholder should understand these controls, even if only some can execute them",
        [
            "Backup and restore tools protect recoverability, but they should only be used for approved operational reasons.",
            "Import and export routines are governed and often produce review or evidence files that should be retained.",
            "Cleanup tools should follow preview or dry-run discipline before any destructive action is allowed.",
            "Storage management, retention windows, orphan-document checks, and notification queue maintenance all affect long-term operational stability.",
            "Training should emphasize that destructive or high-impact actions must be justified, reviewed, and traceable.",
        ],
        font_size=Pt(15.4),
    )
    add_two_column_slide(
        prs,
        "Use Cases",
        next_slide_no(),
        "Core Use Cases: Operational Processing",
        "These scenarios help line users understand where they fit in the bigger process",
        "Front-Line Processing Use Cases",
        [
            "A clerk captures a new staff-due record and verifies intake details",
            "A data entrant imports approved batches and updates claims or registry data",
            "A writeup officer prepares a case and documents what is ready for handoff",
            "A file creator assembles the file package and logs custody movement",
            "A registry or records user updates life-certificate or document detail",
        ],
        "Review and Decision Use Cases",
        [
            "An assessor validates benefits calculations",
            "An auditor checks evidence and control completeness",
            "An approver gives final decision or returns the case",
            "OC/Pension resolves bottlenecks, queues, and high-impact exceptions",
            "An administrator removes access blockers and protects continuity",
        ],
        item_font_size=Pt(14.0),
    )
    add_two_column_slide(
        prs,
        "Use Cases",
        next_slide_no(),
        "Core Use Cases: Service, Governance, and Public Support",
        "The platform also supports service delivery beyond internal line processing",
        "Service and Governance Use Cases",
        [
            "Support staff resolve login, access, and session issues",
            "Payroll and finance staff upload payroll, suspensions, and gratuity schedules",
            "Message users coordinate follow-up across roles",
            "Feedback managers assign, resolve, and export submissions",
            "General users review status, registry, claims, and guidance content",
        ],
        "Pensioner and Public Use Cases",
        [
            "A pensioner checks dashboard status and benefits summary",
            "A pensioner reviews compliance and claim visibility",
            "A pensioner updates allowed profile fields and lookup consent",
            "A pensioner searches the consent-based directory",
            "Public users read guidance pages or submit service feedback",
        ],
        item_font_size=Pt(14.0),
    )
    add_bullet_slide(
        prs,
        "Support",
        next_slide_no(),
        "Troubleshooting and Escalation Checklist",
        "A shared first-response model reduces panic and misdiagnosis",
        [
            "If login fails, check credentials, lockout, maintenance mode, pensioner portal status, and account linkage before escalating.",
            "If an action is missing or denied, confirm the current role, effective role, and permission scope before assuming the feature is broken.",
            "If numbers look wrong, refresh the page, verify filters and date scope, and then confirm in the source module.",
            "If a record appears deleted, check the delete queue or recycle bin before assuming permanent loss.",
            "If an upload or import fails, review template structure, period selection, review exports, and page-specific permissions first.",
            "If a technical or system-health issue remains, escalate with the exact page, action, timestamp, and visible evidence.",
        ],
        font_size=Pt(14.5),
        top=Inches(1.95),
        height=Inches(4.95),
    )
    add_two_column_slide(
        prs,
        "Training Resources",
        next_slide_no(),
        "Training Artefacts and Reference Pack",
        "The platform already has a documentation set that should be used during onboarding and refresher training",
        "Core Reference Documents",
        [
            "System Documentation for current-state platform scope",
            "User and Admin Manual for shared operational practice",
            f"Role manual suite for {inventory['role_manual_count']} maintained roles",
            "Full ERD and domain ERD packs for data understanding",
            f"DFD pack with {inventory['dfd_count']} current process diagrams",
        ],
        "Recommended Training Path",
        [
            "Start with this overview deck for shared context",
            "Move to the relevant role manual for role-specific detail",
            "Run guided walkthroughs in the live or training environment",
            "Use real use cases to practise escalation and handoff",
            "Keep docs, version markers, and procedures aligned after changes",
        ],
        item_font_size=Pt(14.0),
    )
    add_bullet_slide(
        prs,
        "Close",
        next_slide_no(),
        "Closing Guidance and Discussion Prompts",
        "How to use this training effectively after the session",
        [
            "Train by role and scenario, not by showing every page to every user.",
            "Use supervised walk-throughs for high-risk flows such as delete requests, payroll uploads, claims posting, and backup or cleanup operations.",
            "Encourage early escalation whenever role ownership, data state, or system behavior is unclear.",
            "Keep the documentation suite and training deck in the rebuild pipeline so training stays current with the live app.",
            "Discussion prompts: which stakeholder groups need hands-on labs next, which workflows need deeper rehearsal, and which support issues recur most often?",
        ],
        font_size=Pt(15.2),
    )
    return prs


def save_deck(prs: Presentation, path: Path):
    prs.save(str(path))


if __name__ == "__main__":
    save_deck(build_technical_deck(), DOCS / "PensionApp_Technical_Team_Presentation.pptx")
    save_deck(build_user_deck(), DOCS / "PensionApp_Users_Presentation.pptx")
    save_deck(build_training_deck(), DOCS / "PensionApp_Comprehensive_Training_Presentation.pptx")
